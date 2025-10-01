import os, tempfile, shutil, re
from dotenv import load_dotenv
from fastapi import FastAPI, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from openai import OpenAI
from pptx import Presentation
import PyPDF2
from io import BytesIO

load_dotenv()

# OpenAI 클라이언트 초기화
api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    raise ValueError("OPENAI_API_KEY environment variable is not set")

# ASCII로 인코딩 가능한지 확인
try:
    api_key.encode('ascii')
except UnicodeEncodeError:
    # API 키에 비-ASCII 문자가 있으면 제거
    api_key = api_key.encode('ascii', errors='ignore').decode('ascii').strip()

client = OpenAI(api_key=api_key)

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"]
)

@app.get("/health")
async def health_check():
    return {"status": "ok", "message": "Server is running"}

@app.post("/test-openai")
async def test_openai():
    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": "Say hello"}
            ]
        )
        return {"status": "success", "response": response.choices[0].message.content}
    except Exception as e:
        return {"status": "error", "error": str(e)}

def extract_text_from_pptx(file_content: bytes) -> str:
    """PPTX 파일에서 텍스트 추출"""
    try:
        presentation = Presentation(BytesIO(file_content))
        text_content = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            text_content.append(f"=== 슬라이드 {slide_num} ===")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
            
            text_content.append("")  # 슬라이드 간 구분을 위한 빈 줄
        
        return "\n".join(text_content)
    except Exception as e:
        return f"PPTX 파일 처리 중 오류 발생: {str(e)}"

def extract_text_from_pdf(file_content: bytes) -> str:
    """PDF 파일에서 텍스트 추출"""
    try:
        pdf_reader = PyPDF2.PdfReader(BytesIO(file_content))
        text_content = []
        
        for page_num, page in enumerate(pdf_reader.pages, 1):
            text_content.append(f"=== 페이지 {page_num} ===")
            text_content.append(page.extract_text())
            text_content.append("")  # 페이지 간 구분을 위한 빈 줄
        
        return "\n".join(text_content)
    except Exception as e:
        return f"PDF 파일 처리 중 오류 발생: {str(e)}"

@app.post("/stt")
async def stt(
    file: UploadFile = File(...),
    ui_lang: str = Form("ko")
):
    with tempfile.NamedTemporaryFile(delete=False, suffix=f"_{file.filename}") as tmp:
        shutil.copyfileobj(file.file, tmp)
        audio_path = tmp.name
    tx = client.audio.transcriptions.create(
        model="gpt-4o-transcribe",  # 고품질 모델
        file=open(audio_path, "rb"),
        response_format="json"
    )
    transcript = tx.text
    # 선택 언어로 번역 (안내문구 없이 번역 결과만 반환하도록 프롬프트 명확화)
    if ui_lang == "ko":
        sys = (
            "You are a professional translator into Korean. "
            "Translate the following text into natural Korean. "
            "Do NOT add any explanations, introductions, or extra comments. "
            "Return ONLY the translated content itself."
        )
        tgt = "Korean"
    else:
        sys = (
            "You are a professional translator into English. "
            "Translate the following text into natural English. "
            "Do NOT add any explanations, introductions, or extra comments. "
            "Return ONLY the translated content itself."
        )
        tgt = "English"
    tr = client.chat.completions.create(
        model="gpt-4o",  # 더 강력한 LLM
        messages=[
            {"role":"system","content":sys},
            {"role":"user","content":f"Translate to {tgt}:\n\n{transcript}"}
        ]
    )
    out = tr.choices[0].message.content if tr.choices else transcript
    # 후처리: 안내문구 자동 제거
    out = re.sub(r"^(Sure, )?(here is the translation.*?:\s*)?-*\n*", "", out, flags=re.IGNORECASE)
    return {"language": ui_lang, "transcript": out}

@app.post("/multimodal-upload")
async def multimodal_upload(
    files: list[UploadFile] = File(...),
    ui_lang: str = Form("ko"),
    lecture_title: str = Form("")
):
    """멀티모달 파일 업로드 및 텍스트 추출"""
    print(f"Received {len(files)} files for lecture: {lecture_title}")
    all_texts = []
    
    for file in files:
        print(f"Processing file: {file.filename}")
        file_content = await file.read()
        file_extension = file.filename.split('.')[-1].lower()
        
        if file_extension == 'pptx':
            extracted_text = extract_text_from_pptx(file_content)
            all_texts.append(f"[PPT 파일: {file.filename}]\n{extracted_text}")
        elif file_extension == 'pdf':
            extracted_text = extract_text_from_pdf(file_content)
            all_texts.append(f"[PDF 파일: {file.filename}]\n{extracted_text}")
        elif file_extension in ['mp3', 'wav', 'm4a', 'ogg']:
            # 기존 음성 파일 처리 로직
            with tempfile.NamedTemporaryFile(delete=False, suffix=f"_{file.filename}") as tmp:
                tmp.write(file_content)
                audio_path = tmp.name
            
            tx = client.audio.transcriptions.create(
                model="gpt-4o-transcribe",
                file=open(audio_path, "rb"),
                response_format="json"
            )
            transcript = tx.text
            
            # 언어 번역 처리
            if ui_lang == "ko":
                sys = (
                    "You are a professional translator into Korean. "
                    "Translate the following text into natural Korean. "
                    "Do NOT add any explanations, introductions, or extra comments. "
                    "Return ONLY the translated content itself."
                )
                tgt = "Korean"
            else:
                sys = (
                    "You are a professional translator into English. "
                    "Translate the following text into natural English. "
                    "Do NOT add any explanations, introductions, or extra comments. "
                    "Return ONLY the translated content itself."
                )
                tgt = "English"
            
            tr = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role":"system","content":sys},
                    {"role":"user","content":f"Translate to {tgt}:\n\n{transcript}"}
                ]
            )
            translated_text = tr.choices[0].message.content if tr.choices else transcript
            translated_text = re.sub(r"^(Sure, )?(here is the translation.*?:\s*)?-*\n*", "", translated_text, flags=re.IGNORECASE)
            
            all_texts.append(f"[음성 파일: {file.filename}]\n{translated_text}")
            
            # 임시 파일 삭제
            os.unlink(audio_path)
        else:
            all_texts.append(f"[지원하지 않는 파일 형식: {file.filename}]")
    
    # 모든 텍스트를 결합
    combined_text = "\n\n".join(all_texts)
    print(f"Combined text length: {len(combined_text)}")
    print(f"First 200 chars: {combined_text[:200]}")
    
    return {
        "language": ui_lang,
        "combined_text": combined_text,
        "file_count": len(files),
        "lecture_title": lecture_title
    }

@app.post("/summarize")
async def summarize(
    text: str = Form(...),
    target_lang: str = Form("ko"),
    lecture_title: str = Form("")
):
    print(f"Summarizing lecture: {lecture_title}, lang: {target_lang}, text length: {len(text)}")
    
    # 영어로만 프롬프트 작성 (Unicode 인코딩 문제 방지)
    if target_lang == "en":
        sys = (
            f"This lecture is about {lecture_title}. "
            "You are an expert at organizing lecture transcripts. Do not omit important content. "
            "Organize the text into clear, structured sections with headings (e.g., Main Topic, Key Points, Evidence, Examples, Conclusion, Action Items). "
            "Use bullet points, numbering, or tables if appropriate. Do not summarize by shortening, but by structuring and clarifying."
        )
        user_prompt = f"[Lecture Content]\n{text}"
    else:
        sys = (
            f"You are an expert at organizing Korean lecture transcripts about {lecture_title}. "
            "Do not omit important content. Organize the text into clear, structured sections with Korean headings. "
            "Use bullet points, numbering, or tables if appropriate. Do not summarize by shortening, but by structuring and clarifying. "
            "Output must be in Korean."
        )
        user_prompt = f"Please organize the following Korean lecture content:\n\n{text}"
    
    try:
        resp = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role":"system","content":sys},
                {"role":"user","content":user_prompt}
            ]
        )
        out = resp.choices[0].message.content if resp.choices else ""
        print(f"Summary generated, length: {len(out)}")
        return {"summary": out}
    except Exception as e:
        print(f"Error in summarize: {e}")
        raise

@app.post("/quiz")
async def quiz(
    text: str = Form(...),
    target_lang: str = Form("ko"),
    lecture_title: str = Form("")
):
    print(f"Creating quiz for lecture: {lecture_title}, lang: {target_lang}")
    
    # 영어로만 프롬프트 작성 (Unicode 인코딩 문제 방지)
    if target_lang == "en":
        sys = (
            f"You are creating multiple choice questions for a lecture about {lecture_title}. "
            "Create exactly 5 questions with 4 options each (A, B, C, D). "
            "CRITICAL: You must return ONLY a valid JSON array. No other text, no explanations, no markdown formatting. "
            "The response must start with '[' and end with ']'. "
            "Use this exact structure:\n"
            "[\n"
            '  {"id": 1, "question": "Question text?", "options": {"A": "Option A", "B": "Option B", "C": "Option C", "D": "Option D"}, "correct": "A", "explanation": "Explanation why A is correct"},\n'
            '  {"id": 2, "question": "Question text?", "options": {"A": "Option A", "B": "Option B", "C": "Option C", "D": "Option D"}, "correct": "B", "explanation": "Explanation why B is correct"}\n'
            "]\n"
            "Remember: ONLY JSON array, nothing else."
        )
        user_prompt = f"[Lecture Content]\n{text}"
    else:
        sys = (
            f"You are creating Korean multiple choice questions for a lecture about {lecture_title}. "
            "Create exactly 5 questions in Korean with 4 options each (A, B, C, D). "
            "CRITICAL: You must return ONLY a valid JSON array. No other text, no explanations, no markdown formatting. "
            "The response must start with '[' and end with ']'. "
            "Use this exact structure (in Korean):\n"
            "[\n"
            '  {"id": 1, "question": "Korean question text?", "options": {"A": "Korean option A", "B": "Korean option B", "C": "Korean option C", "D": "Korean option D"}, "correct": "A", "explanation": "Korean explanation why A is correct"},\n'
            '  {"id": 2, "question": "Korean question text?", "options": {"A": "Korean option A", "B": "Korean option B", "C": "Korean option C", "D": "Korean option D"}, "correct": "B", "explanation": "Korean explanation why B is correct"}\n'
            "]\n"
            "Remember: ONLY JSON array, nothing else."
        )
        user_prompt = f"Create Korean quiz questions for this Korean lecture content:\n\n{text}"
    
    try:
        resp = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role":"system","content":sys},
                {"role":"user","content":user_prompt}
            ]
        )
    except Exception as e:
        print(f"Error calling OpenAI API: {e}")
        raise
    out = resp.choices[0].message.content if resp.choices else ""
    
    # 디버깅을 위한 로그 출력
    print(f"GPT API Response: {out[:200]}...")  # 처음 200자만 출력
    
    # 응답 후처리: JSON 배열만 추출
    try:
        # JSON 배열 시작과 끝 찾기
        start_idx = out.find('[')
        end_idx = out.rfind(']')
        
        if start_idx != -1 and end_idx != -1 and end_idx > start_idx:
            json_str = out[start_idx:end_idx + 1]
            print(f"Extracted JSON: {json_str[:200]}...")
        else:
            json_str = out
            print("No JSON brackets found, using full response")
        
        # JSON 파싱 시도
        import json
        questions = json.loads(json_str)
        print(f"Successfully parsed {len(questions)} questions")
        return {"questions": questions}
    except json.JSONDecodeError as e:
        print(f"JSON parsing failed: {e}")
        print(f"Full response: {out}")
        
        # JSON 파싱 실패 시 빈 배열 반환 (예시 문제 없음)
        print("Returning empty questions array - no default questions")
        return {"questions": []}

@app.post("/submit-answer")
async def submit_answer(
    question_id: int = Form(...),
    selected_answer: str = Form(...),
    correct_answer: str = Form(...),
    explanation: str = Form(...),
    target_lang: str = Form("ko")
):
    # 답안 제출 결과 및 해설 반환
    is_correct = selected_answer == correct_answer
    
    if target_lang == "ko":
        result_message = "정답입니다!" if is_correct else "틀렸습니다."
    else:
        result_message = "Correct!" if is_correct else "Incorrect."
    
    return {
        "is_correct": is_correct,
        "result_message": result_message,
        "explanation": explanation,
        "correct_answer": correct_answer
    }
